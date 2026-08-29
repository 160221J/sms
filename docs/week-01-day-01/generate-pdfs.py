#!/usr/bin/env python3
"""Render Week 1 Day 1 teaching docs to PDF via headless Chrome + institute branding."""

from __future__ import annotations

import shutil
import subprocess
import sys
import time
from io import BytesIO
from pathlib import Path

import markdown
from pypdf import PdfReader, PdfWriter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
SLIDES = REPO / "docs" / "slides"
PDF_DIR = ROOT / "pdf"
TMP = ROOT / ".pdf-html"
CHROME = "google-chrome"
LOGO = REPO / "docs" / "branding" / "epic-learn-logo.png"
FOOTER = "Epic Learn Institute of Higher Education - Go Programming Master Course"

MD_CSS = """
@page { size: A4; margin: 24mm 16mm 20mm 16mm; }
body {
  font-family: "Segoe UI", Calibri, Liberation Sans, sans-serif;
  color: #14202b;
  line-height: 1.45;
  max-width: 820px;
  margin: 0 auto;
  font-size: 12.5pt;
}
h1 { font-size: 22pt; color: #0a2a5c; margin: 0 0 8pt; }
h2 { font-size: 14pt; color: #0a2a5c; margin: 18pt 0 8pt; }
h3 { font-size: 12.5pt; margin: 14pt 0 6pt; }
p, li { margin: 0 0 6pt; }
pre, code {
  font-family: Consolas, "Liberation Mono", monospace;
  font-size: 10pt;
}
pre {
  background: #10212c;
  color: #e8f2ef;
  padding: 10pt 12pt;
  border-radius: 8pt;
  white-space: pre-wrap;
}
code { background: #eef6f5; padding: 0 3pt; border-radius: 3pt; }
pre code { background: none; color: inherit; padding: 0; }
table { border-collapse: collapse; width: 100%; margin: 8pt 0 14pt; font-size: 11pt; }
th, td { border: 1px solid #d8cfc2; padding: 6pt 8pt; text-align: left; vertical-align: top; }
th { background: #eef6f5; }
hr { border: 0; border-top: 1px solid #d8cfc2; margin: 16pt 0; }
a { color: #0b6e68; }
blockquote { border-left: 4px solid #0b6e68; margin: 8pt 0; padding: 4pt 12pt; color: #4d5d6b; }
"""


def md_to_html(src: Path, title: str) -> Path:
    text = src.read_text(encoding="utf-8")
    body = markdown.markdown(
        text,
        extensions=["tables", "fenced_code", "nl2br"],
    )
    html = f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="utf-8"/>
<title>{title}</title>
<style>{MD_CSS}</style>
</head><body>
{body}
</body></html>
"""
    TMP.mkdir(parents=True, exist_ok=True)
    out = TMP / (src.stem + ".html")
    out.write_text(html, encoding="utf-8")
    return out


def chrome_pdf(src_url: str, dest: Path) -> None:
    dest.parent.mkdir(parents=True, exist_ok=True)
    profile = Path(f"/tmp/chrome-pdf-profile-{dest.stem}")
    shutil.rmtree(profile, ignore_errors=True)
    profile.mkdir(parents=True, exist_ok=True)
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        f"--user-data-dir={profile}",
        "--remote-debugging-port=0",
        "--no-pdf-header-footer",
        "--hide-scrollbars",
        f"--print-to-pdf={dest}",
        src_url,
    ]
    print("PDF", dest.name, flush=True)
    if dest.exists():
        dest.unlink()
    proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    deadline = time.time() + 45
    last_size = -1
    stable = 0
    while time.time() < deadline:
        if dest.exists():
            size = dest.stat().st_size
            if size >= 1000 and size == last_size:
                stable += 1
                if stable >= 3:
                    proc.kill()
                    proc.wait(timeout=5)
                    print(f"  {size // 1024} KB", flush=True)
                    return
            else:
                stable = 0
            last_size = size
        elif proc.poll() is not None:
            break
        time.sleep(0.25)
    proc.kill()
    try:
        proc.wait(timeout=5)
    except subprocess.TimeoutExpired:
        proc.terminate()
    if dest.exists() and dest.stat().st_size >= 1000:
        print(f"  {dest.stat().st_size // 1024} KB", flush=True)
        return
    log = proc.communicate()[0]
    if log:
        sys.stderr.write(log.decode("utf-8", errors="replace"))
    raise SystemExit(f"chrome failed for {dest.name}")


def stamp_branding(pdf_path: Path) -> None:
    """Draw logo header + institute footer on every page."""
    reader = PdfReader(str(pdf_path))
    writer = PdfWriter()
    logo = ImageReader(str(LOGO))
    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        buf = BytesIO()
        c = canvas.Canvas(buf, pagesize=(width, height))
        logo_h = 38
        logo_w = logo_h * (938 / 370)
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, height - 56, width, 56, fill=1, stroke=0)
        c.drawImage(
            logo,
            24,
            height - 50,
            width=logo_w,
            height=logo_h,
            preserveAspectRatio=True,
            mask="auto",
            anchor="sw",
        )
        c.setStrokeColorRGB(0.85, 0.72, 0.15)
        c.setLineWidth(1.2)
        c.line(24, height - 56, width - 24, height - 56)
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, width, 28, fill=1, stroke=0)
        c.setFillColorRGB(0.04, 0.16, 0.36)
        c.setFont("Helvetica", 8)
        c.drawCentredString(width / 2, 12, FOOTER)
        c.setStrokeColorRGB(0.04, 0.16, 0.36)
        c.setLineWidth(0.6)
        c.line(24, 26, width - 24, 26)
        c.save()
        buf.seek(0)
        overlay = PdfReader(buf).pages[0]
        page.merge_page(overlay)
        writer.add_page(page)
    tmp = pdf_path.with_suffix(".stamped.pdf")
    with tmp.open("wb") as fh:
        writer.write(fh)
    tmp.replace(pdf_path)
    print("  branded", pdf_path.name, flush=True)


def pdf_to_pptx(pdf_path: Path, pptx_path: Path) -> None:
    """One PowerPoint slide per PDF page (same look as the branded deck)."""
    import tempfile

    import pymupdf
    from pptx import Presentation
    from pptx.util import Emu, Inches

    width_in, height_in = 13.333, 7.5
    prs = Presentation()
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)
    blank = prs.slide_layouts[6]
    doc = pymupdf.open(pdf_path)
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2), alpha=False)
            img = tmp_path / f"slide-{i:02d}.png"
            pix.save(str(img))
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                str(img),
                Emu(0),
                Emu(0),
                width=prs.slide_width,
                height=prs.slide_height,
            )
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    print("PPTX", pptx_path.name, f"{pptx_path.stat().st_size // 1024} KB", flush=True)


def main() -> int:
    if not LOGO.exists():
        raise SystemExit(f"missing logo: {LOGO}")
    PDF_DIR.mkdir(parents=True, exist_ok=True)
    jobs: list[tuple[str, Path]] = []

    jobs.append((SLIDES.joinpath("week-01-day-01.html").as_uri(), PDF_DIR / "week-01-day-01-slides.pdf"))
    jobs.append((ROOT.joinpath("student-handout.html").as_uri(), PDF_DIR / "student-handout.pdf"))
    jobs.append((ROOT.joinpath("lab-sheet.html").as_uri(), PDF_DIR / "lab-sheet.pdf"))

    md_jobs = [
        (SLIDES / "week-01-day-01-notes.md", "Presenter notes"),
        (ROOT / "pre-class.md", "Pre-class message"),
        (ROOT / "instructor-checklist.md", "Instructor checklist"),
        (ROOT / "install.md", "Install Go, Git, VS Code"),
        (ROOT / "troubleshooting.md", "Troubleshooting"),
        (ROOT / "homework.md", "Homework"),
        (ROOT / "why-go.template.md", "why-go.md template"),
        (ROOT / "README.md", "Day 1 kit index"),
        (REPO / "COURSE.md", "16-week curriculum"),
        (REPO / "docs" / "sms-project-spec.md", "SMS project spec"),
    ]
    for path, title in md_jobs:
        html_path = md_to_html(path, title)
        jobs.append((html_path.as_uri(), PDF_DIR / f"{path.stem}.pdf"))

    for url, dest in jobs:
        chrome_pdf(url, dest)
        stamp_branding(dest)

    slides_pdf = PDF_DIR / "week-01-day-01-slides.pdf"
    pdf_to_pptx(slides_pdf, PDF_DIR / "week-01-day-01-slides.pptx")

    print("Wrote", PDF_DIR)
    for p in sorted(PDF_DIR.glob("*.pdf")):
        print(f"  {p.name:40} {p.stat().st_size // 1024:5d} KB")
    return 0


if __name__ == "__main__":
    sys.exit(main())
