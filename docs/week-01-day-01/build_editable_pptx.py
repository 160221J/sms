#!/usr/bin/env python3
"""Build an editable PowerPoint of the Week 1 Day 1 deck (native text, not screenshots)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parents[1]
LOGO = REPO / "docs" / "branding" / "epic-learn-logo.png"
OUT = ROOT / "pdf" / "week-01-day-01-slides.pptx"
FOOTER = "Epic Learn Institute of Higher Education - Go Programming Master Course"

INK = RGBColor(0x14, 0x20, 0x2B)
MUTED = RGBColor(0x4D, 0x5D, 0x6B)
TEAL = RGBColor(0x0B, 0x6E, 0x68)
TEAL_DARK = RGBColor(0x08, 0x4F, 0x4B)
NAVY = RGBColor(0x0A, 0x2A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF4, 0xEF, 0xE6)
CARD = RGBColor(0xFF, 0xFD, 0xF9)
LINE = RGBColor(0xD8, 0xCF, 0xC2)
CODE_BG = RGBColor(0x10, 0x21, 0x2C)
CODE_FG = RGBColor(0xE8, 0xF2, 0xEF)
GOLD = RGBColor(0xD9, 0xB8, 0x26)
TITLE_BG = RGBColor(0x14, 0x2A, 0x32)
CREAM_TEXT = RGBColor(0xEA, 0xD9, 0xA3)
LIVE = RGBColor(0xC2, 0x4E, 0x1D)

W, H = 13.333, 7.5
MX = 0.65
HEADER_H = 0.78
FOOTER_H = 0.42
CONTENT_TOP = 0.95
CONTENT_BOTTOM = 6.95
TOTAL = 25


def _font(run, name="Calibri", size=18, bold=False, color=INK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_textbox(slide, l, t, w, h, text, size=18, bold=False, color=INK, font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _font(run, font, size, bold, color)
    return box


def add_bullets(slide, l, t, w, h, items, size=20, color=INK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.level = 0
        run = p.add_run()
        run.text = "•  " + item
        _font(run, "Calibri", size, False, color)
    return box


def add_code(slide, l, t, w, h, text, size=15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)
    lines = text.strip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line if line else " "
        _font(run, "Consolas", size, False, CODE_FG)
    return shape


def add_card(slide, l, t, w, h, title, body, title_size=16, body_size=13):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _font(run, "Calibri", title_size, True, TEAL_DARK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    run2 = p2.add_run()
    run2.text = body
    _font(run2, "Calibri", body_size, False, INK)
    return shape


def add_brand(slide, page, dark=False):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(HEADER_H))
    _fill(header, WHITE)
    slide.shapes.add_picture(str(LOGO), Inches(0.35), Inches(0.12), height=Inches(0.55))
    add_textbox(
        slide, 10.2, 0.22, 2.8, 0.4,
        "WEEK 1  ·  DAY 1",
        size=12, bold=True, color=CREAM_TEXT if dark else TEAL_DARK, align=PP_ALIGN.RIGHT,
    )
    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(HEADER_H), Inches(W - 0.7), Inches(0.02)
    )
    _fill(gold, GOLD)

    footer_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(H - FOOTER_H), Inches(W), Inches(FOOTER_H)
    )
    _fill(footer_bg, WHITE)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(H - FOOTER_H), Inches(W - 0.7), Inches(0.015)
    )
    _fill(line, NAVY)
    add_textbox(
        slide, 0.5, H - FOOTER_H + 0.06, 10.5, 0.3,
        FOOTER, size=11, color=NAVY, align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide, 11.3, H - FOOTER_H + 0.06, 1.6, 0.3,
        f"{page} / {TOTAL}", size=11, color=MUTED, align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_kicker_title(slide, kicker, title, title_size=32, top=None):
    top = CONTENT_TOP if top is None else top
    add_textbox(slide, MX, top, 12, 0.32, kicker.upper(), size=13, bold=True, color=TEAL)
    add_textbox(slide, MX, top + 0.28, 12, 0.7, title, size=title_size, bold=True, color=INK)


def new_slide(prs, page, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(H))
        _fill(bg, TITLE_BG)
    else:
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(W), Inches(H))
        _fill(bg, CREAM)
    add_brand(slide, page, dark=dark)
    return slide


def build(dest: Path | None = None) -> Path:
    dest = dest or OUT
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    page = 0

    def P():
        nonlocal page
        page += 1
        return page

    # 1 title
    s = new_slide(prs, P(), dark=True)
    add_textbox(s, MX, 2.15, 12, 0.4, "GO PROGRAMMING MASTER COURSE", size=14, bold=True, color=CREAM_TEXT)
    add_textbox(s, MX, 2.55, 12, 1.0, "Week 1, Day 1", size=48, bold=True, color=WHITE)
    add_textbox(
        s, MX, 3.7, 11, 1.1,
        "Software engineering, why Go, your toolchain, and the first program you will put on GitHub.",
        size=22, color=RGBColor(0xE7, 0xDD, 0xD0),
    )
    add_textbox(
        s, MX, 5.3, 11, 0.9,
        "Umesh Indrajith · BSc Eng (Hons), University of Moratuwa\nSenior Software Engineer / Lecturer",
        size=16, color=RGBColor(0xD7, 0xCB, 0xB8),
    )
    add_notes(s, "Welcome. This is software engineering using Go. SMS starts week 4–5; today is foundations.")

    # 2 who
    s = new_slide(prs, P())
    add_kicker_title(s, "Who is teaching", "Umesh Indrajith")
    cards = [
        ("Background", "BSc Engineering (Hons) in Electronic and Telecommunication Engineering, University of Moratuwa."),
        ("Practice", "Senior Software Engineer. You will learn how Go is used to ship backend systems, not only how the language looks."),
        ("This room", "Ask questions as we go. If the install fails, say so immediately — do not sit broken until the lab."),
    ]
    cw = 3.85
    for i, (t, b) in enumerate(cards):
        add_card(s, MX + i * (cw + 0.18), 2.15, cw, 4.4, t, b, 18, 15)
    add_notes(s, "Keep this short. Credibility, then class norms.")

    # 3 what this course is
    s = new_slide(prs, P())
    add_kicker_title(s, "What this course is", "Software engineering, using Go")
    add_bullets(s, MX, 2.15, 12, 4.4, [
        "Not “memorise syntax for 16 weeks, then a project at the end”.",
        "You will learn to design, build, test, review, and explain working software.",
        "Go is the language. Git, HTTP, databases, and tests are the job.",
        "By week 16 you have a final project in your portfolio: the Student Management System.",
    ], 22)
    add_notes(s, "This is an SE course that uses Go.")

    # 4 shape
    s = new_slide(prs, P())
    add_kicker_title(s, "16 weeks", "The shape of the course")
    cards = [
        ("Weeks 1–3", "Language enough to be dangerous: types, control flow, slices, maps, first functions, Git every week."),
        ("Weeks 4–16", "One product: SMS. Students, then Postgres, then login. Ten to twelve weeks of hands-on."),
        ("Why REST so early?", "So you have something to demo from week 5. We deepen functions and structs inside that project."),
    ]
    for i, (t, b) in enumerate(cards):
        add_card(s, MX + i * (cw + 0.18), 2.15, cw, 4.4, t, b, 18, 15)
    add_notes(s, "Students should leave knowing SMS starts week 4–5.")

    # 5 outcomes
    s = new_slide(prs, P())
    add_kicker_title(s, "Today", "If Day 1 works, you can")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "Explain what software engineering is, in two sentences.",
        "Say why we chose Go for a backend course.",
        "Run Go 1.25, VS Code, and gofmt on Ubuntu (or WSL).",
        "Write, run, and build Hello, Epic Learn.",
        "Put that program on GitHub.",
    ], 22)
    add_notes(s, "SMS is not today.")

    # 6 SE
    s = new_slide(prs, P())
    add_kicker_title(s, "Software engineering", "Building software that other people can trust", 28)
    add_textbox(
        s, MX, 2.1, 12, 1.3,
        "Programming is writing instructions. Software engineering is delivering a system that can be changed, tested, and run by a team — including future you.",
        size=20, color=MUTED,
    )
    add_bullets(s, MX, 3.5, 12, 3.0, [
        "It has users, constraints, and a lifespan longer than the demo.",
        "It has quality: correctness, readability, security of secrets, tests.",
        "It has a process: we do not only “code until it works on my machine”.",
    ], 20)
    add_notes(s, "Ask: who has shipped something another person had to run?")

    # 7 loop
    s = new_slide(prs, P())
    add_kicker_title(s, "The loop we will use all 16 weeks", "Requirements → design → implement → test → review", 24)
    loop = [
        ("Requirements", "Who is it for? What must it do? What is out of scope?"),
        ("Design", "Folders, API shape, data. A little design beats a mess."),
        ("Implement", "Go code. Small steps. Commit often."),
        ("Test & review", "go test, then another human reads it. That is the job."),
    ]
    cw2, ch2 = 5.85, 1.95
    for i, (t, b) in enumerate(loop):
        col, row = i % 2, i // 2
        add_card(s, MX + col * (cw2 + 0.2), 2.15 + row * (ch2 + 0.18), cw2, ch2, t, b, 18, 15)
    add_notes(s, "Draw this on the board and leave it up for the term.")

    # 8 how we work
    s = new_slide(prs, P())
    add_kicker_title(s, "A course, not a playlist", "How we work here")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "I type, you type. Watching is not the skill.",
        "Homework is a GitHub pull request, not a zip on WhatsApp.",
        "Unformatted code comes back. gofmt is not optional.",
        "AI tools are allowed. If you cannot explain the line, it is not yours.",
    ], 22)
    add_notes(s, "Be warm but firm on the AI policy.")

    # 9 why go
    s = new_slide(prs, P())
    add_kicker_title(s, "Why Go", "A language designed for building services", 28)
    add_textbox(
        s, MX, 2.1, 12, 1.15,
        "Created at Google (2007–2009) for large systems: simple language, excellent toolchain, great concurrency, one static binary.",
        size=18, color=MUTED,
    )
    add_bullets(s, MX, 3.35, 12, 3.2, [
        "Used for APIs, CLIs, cloud tooling, Kubernetes-related systems.",
        "Fast to compile. Easy to read. Hard to hide complexity behind magic.",
        "We chose it because a beginner can reach a real HTTP API in this course.",
    ], 20)
    add_notes(s, "Do not bash Java/Python.")

    # 10 features
    s = new_slide(prs, P())
    add_kicker_title(s, "Why Go · features you will actually use", "What we are buying")
    feats = [
        ("Simplicity", "Small language. Fewer ways to write the same thing. Good for teams."),
        ("Toolchain", "go run, go build, gofmt, go test, go mod — one installer."),
        ("Concurrency", "Goroutines. Weeks 10–12. Mention now, practise later."),
        ("One binary", "Compile and copy. No virtualenv ritual to run the server."),
        ("Standard library", "HTTP and JSON live in stdlib. We add Gin once you have seen net/http."),
        ("Garbage collection", "You will still learn pointers. You will not manage malloc by hand."),
    ]
    cw3, ch3 = 3.85, 2.15
    for i, (t, b) in enumerate(feats):
        col, row = i % 3, i // 3
        add_card(s, MX + col * (cw3 + 0.18), 2.05 + row * (ch3 + 0.16), cw3, ch3, t, b, 16, 13)
    add_notes(s, "Pause on toolchain — that is today's practical why.")

    # 11 SMS
    s = new_slide(prs, P())
    add_kicker_title(s, "Where we are going", "Student Management System (SMS)")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "Staff register and log in.",
        "Staff create, list, update, delete students.",
        "Go API + PostgreSQL. A React UI is provided — you are graded on Go.",
        "Not today. Weeks 1–3 are the tools. Project kickoff is week 4.",
    ], 22)
    add_notes(s, "Do not dive into Gin today.")

    # 12 tools
    s = new_slide(prs, P())
    add_kicker_title(s, "Pinned tools", "Everyone uses the same versions")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "Go 1.25.x from go.dev — not whatever apt install golang-go gives you.",
        "Ubuntu in the lab. Windows: WSL2 + Ubuntu. macOS is fine.",
        "VS Code + official Go extension. GoLand is optional if you already have it.",
        "Git + a GitHub account. Postman comes when we do HTTP.",
    ], 20)
    add_notes(s, "Confirm nobody uses an ancient apt Go.")

    # 13 install demo
    s = new_slide(prs, P())
    add_textbox(s, MX, CONTENT_TOP, 2.2, 0.32, "LIVE DEMO", size=12, bold=True, color=LIVE)
    add_kicker_title(s, "Install Go 1.25 on Ubuntu", "Do this with me", top=1.22)
    add_code(s, MX, 2.35, 12, 3.5, """wget https://go.dev/dl/go1.25.0.linux-amd64.tar.gz
sudo rm -rf /usr/local/go
sudo tar -C /usr/local -xzf go1.25.0.linux-amd64.tar.gz
echo 'export PATH=$PATH:/usr/local/go/bin' >> ~/.bashrc
source ~/.bashrc
go version""", 16)
    add_notes(s, "Use the exact patch version on go.dev/dl. Do not apt install golang-go.")

    # 14 go env
    s = new_slide(prs, P())
    add_kicker_title(s, "Did it work?", "go version and go env")
    add_code(s, MX, 2.1, 12, 1.35, "go version\ngo env GOROOT GOPATH GOMODCACHE GO111MODULE", 16)
    add_bullets(s, MX, 3.6, 12, 3.0, [
        "GOROOT — where the Go install lives (usually /usr/local/go).",
        "GOPATH — old workspace idea. We do not put class projects there.",
        "Modules — a folder with go.mod is the project. That is week 4. Today: one file is enough.",
    ], 18)
    add_notes(s, "Ignore GOPATH for this course except that it exists.")

    # 15 vscode
    s = new_slide(prs, P())
    add_textbox(s, MX, CONTENT_TOP, 2.2, 0.32, "LIVE DEMO", size=12, bold=True, color=LIVE)
    add_kicker_title(s, "IDE", "VS Code, set up for Go", top=1.22)
    add_bullets(s, MX, 2.3, 12, 4.3, [
        "Install VS Code if needed. Open it.",
        "Extensions: search Go by the Go Team at Google. Install it.",
        "Open a folder you will use for this course, e.g. ~/epic-go.",
        "Command Palette: Go: Install/Update Tools — install all (gofmt, gopls, dlv).",
    ], 20)
    add_notes(s, "VS Code is the class default.")

    # 16 format on save
    s = new_slide(prs, P())
    add_kicker_title(s, "The setting you must turn on", "Format on save = gofmt")
    add_code(s, MX, 2.1, 12, 2.5, """{
  "editor.formatOnSave": true,
  "[go]": {
    "editor.defaultFormatter": "golang.go",
    "editor.formatOnSave": true
  }
}""", 16)
    add_textbox(
        s, MX, 4.8, 12, 1.4,
        "Settings JSON (Ctrl+Shift+P → Preferences: Open User Settings (JSON)). From today, unformatted Go is rejected.",
        size=18, color=MUTED,
    )
    add_notes(s, "Show before/after: messy hello.go, save, it snaps into shape.")

    # 17 hello.go
    s = new_slide(prs, P())
    add_kicker_title(s, "First program", "Anatomy of a Go file")
    add_code(s, MX, 2.05, 12, 2.35, """package main

import "fmt"

func main() {
    fmt.Println("Hello, Epic Learn")
}""", 16)
    add_bullets(s, MX, 4.55, 12, 2.1, [
        "package main + func main → an executable program.",
        'import "fmt" — standard library formatted I/O.',
        "Save as hello.go.",
    ], 18)
    add_notes(s, "Type it from scratch. Students type with you.")

    # 18 export
    s = new_slide(prs, P())
    add_kicker_title(s, "Names in Go", "Exported vs unexported")
    add_bullets(s, MX, 2.1, 6.6, 4.5, [
        "A name that starts with a capital letter is exported (visible from other packages).",
        "fmt.Println works because Println is exported.",
        "fmt.println will not compile.",
        "This is not Java’s public keyword. The capital letter is the rule.",
    ], 18)
    add_code(s, 7.5, 2.2, 5.15, 2.4, 'fmt.Println("ok")\n\n// fmt.println("no")\n// undefined: fmt.println', 16)
    add_notes(s, "Do not use Main vs main as the example. Println is the clean example.")

    # 19 three commands
    s = new_slide(prs, P())
    add_textbox(s, MX, CONTENT_TOP, 2.2, 0.32, "LIVE DEMO", size=12, bold=True, color=LIVE)
    add_kicker_title(s, "Develop, compile, run", "Three commands", top=1.22)
    add_code(s, MX, 2.25, 12, 2.25, """gofmt -w hello.go           # rewrite the file in standard format
go run hello.go             # compile in a temp place and run
go build -o hello hello.go  # produce a binary (no go.mod needed yet)
./hello                     # run the binary""", 15)
    add_bullets(s, MX, 4.6, 12, 2.0, [
        "go run — while learning.",
        "go build — what you ship (a file you can copy).",
        "go test exists. We use it on a real function in week 2.",
    ], 18)
    add_notes(s, "Do not go build . today — there is no go.mod until week 4.")

    # 20 git why
    s = new_slide(prs, P())
    add_kicker_title(s, "Version control", "Why Git on day 1")
    add_textbox(
        s, MX, 2.4, 12, 3.5,
        "If it is not on GitHub, it is a story you told. The course artefact is the history of your work: commits, then pull requests from week 4.",
        size=28, bold=True, color=INK,
    )
    add_notes(s, "Make GitHub a completion gate for the lab.")

    # 21 git commands
    s = new_slide(prs, P())
    add_textbox(s, MX, CONTENT_TOP, 2.2, 0.32, "LIVE DEMO", size=12, bold=True, color=LIVE)
    add_kicker_title(s, "Git", "Init, ignore, commit", top=1.22)
    add_code(s, MX, 2.25, 12, 4.25, """git config --global user.name  "Your Name"
git config --global user.email "you@example.com"

git init
echo hello > .gitignore
git add hello.go .gitignore
git status
git commit -m "Add Hello Epic Learn" """, 15)
    add_notes(s, "The binary from go build must not be committed.")

    # 22 github
    s = new_slide(prs, P())
    add_kicker_title(s, "GitHub", "A remote is a backup and a classroom")
    add_code(s, MX, 2.05, 12, 1.7, """git branch -M main
git remote add origin https://github.com/YOU/hello-epic-learn.git
git push -u origin main""", 15)
    add_bullets(s, MX, 3.95, 12, 2.6, [
        "Create an empty repo on github.com (no README if you already committed locally).",
        "Add me as collaborator if that is how this batch submits.",
        "git clone is how you copy a repo onto a new machine — we will use it for SMS in week 4.",
    ], 18)
    add_notes(s, "Walk the room during push.")

    # 23 lab
    s = new_slide(prs, P())
    add_kicker_title(s, "Lab · rest of class", "Done when all of this is true")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "go version shows 1.25.x",
        "VS Code formats a Go file on save",
        "hello.go prints Hello, Epic Learn via go run and via ./hello",
        "GitHub repo contains hello.go, not the compiled binary",
        "Show me the GitHub URL before you leave",
    ], 20)
    add_notes(s, "Homework is on the student handout — not a slide.")

    # 24 next
    s = new_slide(prs, P())
    add_kicker_title(s, "Next session", "Week 1 wrap / Week 2 start")
    add_bullets(s, MX, 2.15, 12, 4.5, [
        "If anyone’s install is still broken, we fix it first.",
        "Then: variables, types, zero values, operators, strings.",
        "Bring a working go version and your GitHub URL.",
    ], 22)
    add_notes(s, "Do not start slices until week 3.")

    # 25 questions
    s = new_slide(prs, P(), dark=True)
    add_textbox(s, 0.5, 2.9, 12.3, 1.6, "Questions", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_notes(s, "Take questions. Then start the install lab.")

    if page != TOTAL:
        raise SystemExit(f"expected {TOTAL} slides, built {page}")

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    print(f"PPTX {dest} ({dest.stat().st_size // 1024} KB, {page} editable slides)")
    return dest


if __name__ == "__main__":
    build()
